import PyPDF2
import pdfplumber
from docx import Document
from pptx import Presentation
import os


class FileExtractor:
    """Extract text from different file formats"""
    
    @staticmethod
    def extract_from_pdf(file_path):
        """Extract text from PDF file"""
        text = ""
        try:
            # Try with pdfplumber first (better for complex PDFs)
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            
            # If no text found, try PyPDF2
            if not text.strip():
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
            
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting PDF: {str(e)}")
    
    @staticmethod
    def extract_from_docx(file_path):
        """Extract text from DOCX file"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + "\n"
            
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting DOCX: {str(e)}")
    
    @staticmethod
    def extract_from_pptx(file_path):
        """Extract text from PPTX file"""
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_number, slide in enumerate(prs.slides, start=1):
                text += f"\n--- Slide {slide_number} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting PPTX: {str(e)}")
    
    @staticmethod
    def extract_text(file_path):
        """Main method to extract text based on file extension"""
        _, extension = os.path.splitext(file_path)
        extension = extension.lower()
        
        if extension == '.pdf':
            return FileExtractor.extract_from_pdf(file_path)
        elif extension == '.docx':
            return FileExtractor.extract_from_docx(file_path)
        elif extension == '.pptx':
            return FileExtractor.extract_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {extension}")
    
    @staticmethod
    def count_words(text):
        """Count words in extracted text"""
        words = text.split()
        return len(words)
